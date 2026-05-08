from flask import Flask, request, send_file
from flask_cors import CORS
import os, uuid

from PyPDF2 import PdfReader, PdfWriter
from docx import Document
import pdfplumber
from pdf2image import convert_from_path
from reportlab.pdfgen import canvas
import pytesseract
from PIL import Image
import pandas as pd

from pptx import Presentation
from pptx.util import Inches
from io import BytesIO

app = Flask(__name__)
CORS(app)

UPLOAD = "uploads"
OUTPUT = "outputs"

os.makedirs(UPLOAD, exist_ok=True)
os.makedirs(OUTPUT, exist_ok=True)

# ---------------- HOME ----------------
@app.route("/")
def home():
    return "PDFStudio API running 🚀"

# ---------------- PDF TO WORD ----------------
@app.route("/pdf-to-word", methods=["POST"])
def pdf_to_word():
    f = request.files["files"]
    path = os.path.join(UPLOAD, f.filename)
    f.save(path)

    out = os.path.join(OUTPUT, str(uuid.uuid4()) + ".docx")

    doc = Document()
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                doc.add_paragraph(text)

    doc.save(out)
    return send_file(out, as_attachment=True)

# ---------------- WORD TO PDF ----------------
@app.route("/word-to-pdf", methods=["POST"])
def word_to_pdf():
    f = request.files["files"]
    path = os.path.join(UPLOAD, f.filename)
    f.save(path)

    out = os.path.join(OUTPUT, str(uuid.uuid4()) + ".pdf")

    doc = Document(path)
    c = canvas.Canvas(out)

    y = 800
    for para in doc.paragraphs:
        c.drawString(40, y, para.text[:100])
        y -= 15
        if y < 40:
            c.showPage()
            y = 800

    c.save()
    return send_file(out, as_attachment=True)

# ---------------- PDF TO EXCEL ----------------
@app.route("/pdf-to-excel", methods=["POST"])
def pdf_to_excel():
    f = request.files["files"]
    path = os.path.join(UPLOAD, f.filename)
    f.save(path)

    rows = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            table = page.extract_table()
            if table:
                rows += table

    df = pd.DataFrame(rows)
    out = os.path.join(OUTPUT, str(uuid.uuid4()) + ".xlsx")
    df.to_excel(out, index=False)

    return send_file(out, as_attachment=True)

# ---------------- EXCEL TO PDF ----------------
@app.route("/excel-to-pdf", methods=["POST"])
def excel_to_pdf():
    f = request.files["files"]
    path = os.path.join(UPLOAD, f.filename)
    f.save(path)

    df = pd.read_excel(path)
    out = os.path.join(OUTPUT, str(uuid.uuid4()) + ".pdf")

    c = canvas.Canvas(out)
    y = 800

    for row in df.values:
        c.drawString(40, y, " | ".join(map(str, row)))
        y -= 15
        if y < 40:
            c.showPage()
            y = 800

    c.save()
    return send_file(out, as_attachment=True)

# ---------------- PDF TO JPG ----------------
@app.route("/pdf-to-jpg", methods=["POST"])
def pdf_to_jpg():
    f = request.files["files"]
    path = os.path.join(UPLOAD, f.filename)
    f.save(path)

    images = convert_from_path(path)
    out = os.path.join(OUTPUT, "page.jpg")

    images[0].save(out, "JPEG")
    return send_file(out, as_attachment=True)

# ---------------- JPG TO PDF ----------------
@app.route("/jpg-to-pdf", methods=["POST"])
def jpg_to_pdf():
    f = request.files["files"]
    path = os.path.join(UPLOAD, f.filename)
    f.save(path)

    img = Image.open(path)
    out = os.path.join(OUTPUT, str(uuid.uuid4()) + ".pdf")
    img.convert("RGB").save(out)

    return send_file(out, as_attachment=True)

# ---------------- MERGE PDF ----------------
@app.route("/merge", methods=["POST"])
def merge():
    files = request.files.getlist("files")
    writer = PdfWriter()

    out = os.path.join(OUTPUT, "merged.pdf")

    for f in files:
        path = os.path.join(UPLOAD, f.filename)
        f.save(path)

        reader = PdfReader(path)
        for page in reader.pages:
            writer.add_page(page)

    with open(out, "wb") as o:
        writer.write(o)

    return send_file(out, as_attachment=True)

# ---------------- SPLIT PDF ----------------
@app.route("/split", methods=["POST"])
def split():
    f = request.files["files"]
    path = os.path.join(UPLOAD, f.filename)
    f.save(path)

    reader = PdfReader(path)
    writer = PdfWriter()

    writer.add_page(reader.pages[0])

    out = os.path.join(OUTPUT, "split.pdf")

    with open(out, "wb") as o:
        writer.write(o)

    return send_file(out, as_attachment=True)

# ---------------- COMPRESS (BASIC) ----------------
@app.route("/compress", methods=["POST"])
def compress():
    f = request.files["files"]
    path = os.path.join(UPLOAD, f.filename)
    f.save(path)

    reader = PdfReader(path)
    writer = PdfWriter()

    for page in reader.pages:
        writer.add_page(page)

    out = os.path.join(OUTPUT, str(uuid.uuid4()) + ".pdf")

    with open(out, "wb") as o:
        writer.write(o)

    return send_file(out, as_attachment=True)

# ---------------- PROTECT ----------------
@app.route("/protect", methods=["POST"])
def protect():
    f = request.files["files"]
    path = os.path.join(UPLOAD, f.filename)
    f.save(path)

    reader = PdfReader(path)
    writer = PdfWriter()

    for page in reader.pages:
        writer.add_page(page)

    writer.encrypt("1234")

    out = os.path.join(OUTPUT, str(uuid.uuid4()) + ".pdf")

    with open(out, "wb") as o:
        writer.write(o)

    return send_file(out, as_attachment=True)

# ---------------- OCR ----------------
@app.route("/ocr", methods=["POST"])
def ocr():
    f = request.files["files"]
    path = os.path.join(UPLOAD, f.filename)
    f.save(path)

    images = convert_from_path(path)
    out = os.path.join(OUTPUT, str(uuid.uuid4()) + ".pdf")

    c = canvas.Canvas(out)

    for img in images:
        text = pytesseract.image_to_string(img, lang="eng")
        c.drawString(40, 800, text[:200])
        c.showPage()

    c.save()
    return send_file(out, as_attachment=True)

# ---------------- POWERPOINT TO PDF ----------------
@app.route("/powerpoint-to-pdf", methods=["POST"])
def ppt_to_pdf():
    f = request.files["files"]
    path = os.path.join(UPLOAD, f.filename)
    f.save(path)

    prs = Presentation(path)
    out = os.path.join(OUTPUT, str(uuid.uuid4()) + ".pdf")

    c = canvas.Canvas(out)

    for slide in prs.slides:
        text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

        c.drawString(40, 800, text[:200])
        c.showPage()

    c.save()
    return send_file(out, as_attachment=True)

# ---------------- PDF TO POWERPOINT ----------------
@app.route("/pdf-to-powerpoint", methods=["POST"])
def pdf_to_ppt():
    f = request.files["files"]
    path = os.path.join(UPLOAD, f.filename)
    f.save(path)

    prs = Presentation()
    out = os.path.join(OUTPUT, str(uuid.uuid4()) + ".pptx")

    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            text = page.extract_text()

            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Page"
            slide.placeholders[1].text = text if text else ""

    prs.save(out)
    return send_file(out, as_attachment=True)

# ---------------- RUN (RAILWAY SAFE) ----------------
# IMPORTANT: DO NOT use app.run() on Railway

app = app
